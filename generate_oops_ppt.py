"""
generate_oops_ppt.py
--------------------
Generates a 17-slide PowerPoint presentation on the 4 Pillars of
Object-Oriented Programming using the python-pptx library.

Usage:
    python generate_oops_ppt.py

Output:
    OOP_Pillars_Presentation.pptx  (created in the current directory)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
DARK_BLUE   = RGBColor(0x1F, 0x39, 0x7A)   # titles / header rows
MED_BLUE    = RGBColor(0x2E, 0x74, 0xB5)   # accent / subtitle
LIGHT_BLUE  = RGBColor(0xBD, 0xD7, 0xEE)   # alternating table rows
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x26, 0x26, 0x26)   # body text
CODE_BG     = RGBColor(0xF2, 0xF2, 0xF2)   # light-gray code box background
CODE_FG     = RGBColor(0x1E, 0x1E, 0x1E)   # dark text in code box

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helper: set slide background to white
# ---------------------------------------------------------------------------
def set_slide_background(slide):
    """Fill the slide background with solid white."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


# ---------------------------------------------------------------------------
# Helper: add a coloured rectangle (title bar or code box)
# ---------------------------------------------------------------------------
def _add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()   # no border
    return shape


# ---------------------------------------------------------------------------
# Helper: add a slide-number text box (bottom-right)
# ---------------------------------------------------------------------------
def _add_slide_number(slide, number):
    txBox = slide.shapes.add_textbox(
        Inches(12.5), Inches(7.1), Inches(0.6), Inches(0.3)
    )
    tf = txBox.text_frame
    tf.text = str(number)
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(12)
    run.font.color.rgb = MED_BLUE
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT


# ---------------------------------------------------------------------------
# Helper: add a blue title bar + title text
# ---------------------------------------------------------------------------
def _add_title_bar(slide, title_text):
    bar = _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.15), DARK_BLUE)

    txBox = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.12), Inches(12.7), Inches(0.9)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.bold = True
    run.font.size = Pt(30)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"
    return bar


# ---------------------------------------------------------------------------
# 1. Title Slide
# ---------------------------------------------------------------------------
def add_title_slide(prs):
    """Slide 1 — decorative title slide, no slide number."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)

    # Top accent bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.25), MED_BLUE)

    # Large blue centre block
    _add_rect(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(3.6), DARK_BLUE)

    # Bottom accent bar
    _add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), MED_BLUE)

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.85), Inches(12.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Pillars of Object-Oriented Programming"
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = WHITE
    run.font.name = "Calibri"

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.4), Inches(12.3), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = "Understanding the 4 Core Principles of OOP"
    run2.font.size = Pt(22)
    run2.font.color.rgb = LIGHT_BLUE
    run2.font.name = "Calibri"

    # Date
    tb3 = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.5))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = "March 17, 2026"
    run3.font.size = Pt(18)
    run3.font.color.rgb = DARK_BLUE
    run3.font.name = "Calibri"


# ---------------------------------------------------------------------------
# 2. Generic content slide (bullet points)
# ---------------------------------------------------------------------------
def add_content_slide(prs, title, bullets, slide_num):
    """
    Add a slide with a blue title bar and bullet-point body text.

    bullets  : list of (indent_level, text) tuples
               indent_level 0 = main bullet, 1 = sub-bullet
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)
    _add_title_bar(slide, title)
    _add_slide_number(slide, slide_num)

    # Body text box
    tb = slide.shapes.add_textbox(
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.9)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    for i, (level, text) in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        # Indent via spaces for sub-bullets
        indent = "    " * level
        run = p.add_run()
        run.text = indent + text
        run.font.size = Pt(19 if level == 0 else 17)
        run.font.color.rgb = DARK_GRAY
        run.font.name = "Calibri"
        run.font.bold = (level == 0 and text.startswith("•") is False
                         and "::" in text)
        # Add a small space after each paragraph
        p.space_after = Pt(4)


# ---------------------------------------------------------------------------
# 3. Code slide
# ---------------------------------------------------------------------------
def add_code_slide(prs, title, code_lines, caption, slide_num):
    """
    Add a slide with a blue title bar, a light-gray code box (Courier New),
    and an optional caption below.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)
    _add_title_bar(slide, title)
    _add_slide_number(slide, slide_num)

    # Gray background box for code
    box_top    = Inches(1.25)
    box_height = Inches(5.2)
    _add_rect(slide, Inches(0.3), box_top, Inches(12.7), box_height, CODE_BG)

    # Code text box (on top of the gray box)
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.35), Inches(12.5), Inches(5.0)
    )
    tf = tb.text_frame
    tf.word_wrap = False

    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = "Courier New"
        run.font.size = Pt(14)
        run.font.color.rgb = CODE_FG

    # Caption below the code box
    if caption:
        tb2 = slide.shapes.add_textbox(
            Inches(0.4), Inches(6.55), Inches(12.5), Inches(0.65)
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = caption
        run2.font.size = Pt(17)
        run2.font.color.rgb = MED_BLUE
        run2.font.name = "Calibri"
        run2.font.italic = True


# ---------------------------------------------------------------------------
# 4. Table slide
# ---------------------------------------------------------------------------
def add_table_slide(prs, title, headers, rows, slide_num):
    """
    Add a slide with a blue title bar and a formatted comparison table.

    headers : list of column header strings
    rows    : list of row tuples (one tuple per data row)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide)
    _add_title_bar(slide, title)
    _add_slide_number(slide, slide_num)

    cols = len(headers)
    data_rows = len(rows)
    total_rows = data_rows + 1  # +1 for header

    # Determine table dimensions
    tbl_left   = Inches(0.4)
    tbl_top    = Inches(1.3)
    tbl_width  = Inches(12.5)
    tbl_height = Inches(5.8)

    table = slide.shapes.add_table(
        total_rows, cols, tbl_left, tbl_top, tbl_width, tbl_height
    ).table

    # Equal column widths
    col_width = tbl_width // cols
    for col in table.columns:
        col.width = col_width

    # Header row
    for ci, header in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.bold = True
        run.font.size = Pt(17)
        run.font.color.rgb = WHITE
        run.font.name = "Calibri"

    # Data rows (alternating colours)
    for ri, row in enumerate(rows):
        bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(16)
            run.font.color.rgb = DARK_GRAY
            run.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Build all 17 slides
# ---------------------------------------------------------------------------
def build_presentation():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # ------------------------------------------------------------------
    # Slide 1 — Title
    # ------------------------------------------------------------------
    add_title_slide(prs)

    # ------------------------------------------------------------------
    # Slide 2 — What is OOP?
    # ------------------------------------------------------------------
    add_content_slide(prs, "What is OOP?", [
        (0, "OOP is a programming paradigm based on the concept of \"objects\""),
        (1, "Objects contain data (attributes) and behaviour (methods)"),
        (1, "Models real-world entities in code"),
        (1, "Popular OOP languages: Java, C++, Python, C#"),
        (0, "Real-world analogy:"),
        (1, "Everything around you is an object — a car, a student, a phone"),
        (1, "Each has properties (attributes) and actions (methods)"),
        (0, "Why OOP?"),
        (1, "Organised, reusable, and scalable code"),
        (1, "Easier to maintain and extend large applications"),
    ], slide_num=2)

    # ------------------------------------------------------------------
    # Slide 3 — 4 Pillars Overview (table)
    # ------------------------------------------------------------------
    add_table_slide(prs,
        "The 4 Pillars of OOP — Overview",
        ["#", "Pillar", "One-Liner"],
        [
            ("1", "Encapsulation 🔒", "Wrapping data & methods together, restricting direct access"),
            ("2", "Abstraction 🎭",   "Hiding complexity, showing only essential features"),
            ("3", "Inheritance 👨‍👧‍👦", "Reusing properties & methods from a parent class"),
            ("4", "Polymorphism 🔀",  "One name, many forms — same method, different behaviour"),
        ],
        slide_num=3,
    )

    # ------------------------------------------------------------------
    # Slide 4 — Pillar 1: Encapsulation
    # ------------------------------------------------------------------
    add_content_slide(prs, "Pillar 1: Encapsulation 🔒", [
        (0, "Definition:"),
        (1, "Bundling data (variables) and methods (functions) into a single unit (class)"),
        (1, "Restricting direct access to some components of an object"),
        (0, "Real-life examples:"),
        (1, "Capsule 💊 — medicine is packed inside; you just take it, not mix it"),
        (1, "ATM 🏧 — you interact with buttons; the internal cash mechanism is hidden"),
        (0, "Key concepts:"),
        (1, "Use private variables"),
        (1, "Provide public getters/setters to access or modify them"),
        (1, "Protects data from unintended modification"),
    ], slide_num=4)

    # ------------------------------------------------------------------
    # Slide 5 — Encapsulation Code Example (Java)
    # ------------------------------------------------------------------
    add_code_slide(prs,
        "Encapsulation — Code Example (Java)",
        [
            "class Student {",
            "    private String name;  // hidden from outside",
            "    private int age;",
            "",
            "    // Getter",
            "    public String getName() { return name; }",
            "",
            "    // Setter",
            "    public void setName(String name) {",
            "        this.name = name;",
            "    }",
            "}",
        ],
        "✅  Benefit: Data is safe. No one can directly change 'name' without going through setName().",
        slide_num=5,
    )

    # ------------------------------------------------------------------
    # Slide 6 — Pillar 2: Abstraction
    # ------------------------------------------------------------------
    add_content_slide(prs, "Pillar 2: Abstraction 🎭", [
        (0, "Definition:"),
        (1, "Hiding the internal implementation details"),
        (1, "Showing only the essential features to the user"),
        (0, "Real-life examples:"),
        (1, "Driving a car 🚗 — you use the steering wheel & pedals; you don't need to know how the engine works"),
        (1, "TV Remote 📺 — you press buttons; the circuit inside is irrelevant to you"),
        (0, "How to achieve Abstraction in Java:"),
        (1, "Using abstract classes"),
        (1, "Using interfaces"),
        (0, "Key idea: What an object does vs. How it does it"),
    ], slide_num=6)

    # ------------------------------------------------------------------
    # Slide 7 — Abstraction Code Example (Java)
    # ------------------------------------------------------------------
    add_code_slide(prs,
        "Abstraction — Code Example (Java)",
        [
            "abstract class Shape {",
            "    abstract void draw();  // no implementation here",
            "}",
            "",
            "class Circle extends Shape {",
            "    void draw() {",
            "        System.out.println(\"Drawing a Circle\");",
            "    }",
            "}",
            "",
            "class Rectangle extends Shape {",
            "    void draw() {",
            "        System.out.println(\"Drawing a Rectangle\");",
            "    }",
            "}",
        ],
        "✅  Benefit: User only knows draw() exists — doesn't care HOW each shape draws itself.",
        slide_num=7,
    )

    # ------------------------------------------------------------------
    # Slide 8 — Encapsulation vs Abstraction
    # ------------------------------------------------------------------
    add_table_slide(prs,
        "Encapsulation vs Abstraction ⚖️",
        ["Feature", "Encapsulation", "Abstraction"],
        [
            ("Focus",   "Data hiding",                      "Implementation hiding"),
            ("How",     "Access modifiers (private/public)", "Abstract classes & interfaces"),
            ("Purpose", "Protect data",                     "Reduce complexity"),
            ("Analogy", "Capsule (medicine packed inside)", "Car dashboard (shows speed, not engine)"),
        ],
        slide_num=8,
    )

    # ------------------------------------------------------------------
    # Slide 9 — Pillar 3: Inheritance
    # ------------------------------------------------------------------
    add_content_slide(prs, "Pillar 3: Inheritance 👨‍👧‍👦", [
        (0, "Definition:"),
        (1, "A class (child/subclass) inherits properties and methods from another class (parent/superclass)"),
        (1, "Promotes code reusability"),
        (0, "Real-life examples:"),
        (1, "A child inherits traits from parents — eyes, height, etc."),
        (1, "A sports car 🏎️ inherits features of a generic Car but adds turbo mode"),
        (0, "Types of Inheritance:"),
        (1, "Single | Multilevel | Hierarchical | Multiple (via interfaces) | Hybrid"),
        (0, "Java keyword:  extends  (for classes),  implements  (for interfaces)"),
    ], slide_num=9)

    # ------------------------------------------------------------------
    # Slide 10 — Inheritance Code Example (Java)
    # ------------------------------------------------------------------
    add_code_slide(prs,
        "Inheritance — Code Example (Java)",
        [
            "class Animal {",
            "    void eat() {",
            "        System.out.println(\"This animal eats food\");",
            "    }",
            "}",
            "",
            "class Dog extends Animal {",
            "    void bark() {",
            "        System.out.println(\"Dog barks!\");",
            "    }",
            "}",
            "",
            "// Usage:",
            "Dog d = new Dog();",
            "d.eat();   // Inherited from Animal ✅",
            "d.bark();  // Dog's own method   ✅",
        ],
        "✅  Benefit: Code reusability — no need to rewrite eat() inside the Dog class.",
        slide_num=10,
    )

    # ------------------------------------------------------------------
    # Slide 11 — Types of Inheritance Diagram
    # ------------------------------------------------------------------
    add_content_slide(prs, "Types of Inheritance — Diagram", [
        (0, "Single Inheritance:"),
        (1, "[Parent A]  →  [Child B]"),
        (0, "Multilevel Inheritance:"),
        (1, "[A]  →  [B]  →  [C]         (grandchild inherits from grandparent)"),
        (0, "Hierarchical Inheritance:"),
        (1, "              [A]"),
        (1, "            ↙       ↘"),
        (1, "          [B]       [C]      (two children share one parent)"),
        (0, "Multiple Inheritance (via Interfaces in Java):"),
        (1, "[InterfaceX] + [InterfaceY]  →  [Class Z  implements X, Y]"),
        (0, "💡 Java does NOT support multiple class-inheritance to avoid the Diamond Problem."),
    ], slide_num=11)

    # ------------------------------------------------------------------
    # Slide 12 — Pillar 4: Polymorphism
    # ------------------------------------------------------------------
    add_content_slide(prs, "Pillar 4: Polymorphism 🎭🔀", [
        (0, "Definition:"),
        (1, "The ability of an object to take many forms"),
        (1, "Same method name, different behaviour"),
        (0, "Real-life examples:"),
        (1, "A person 👤 can be a student, an employee, and a parent — same person, different roles"),
        (1, "The \"+\" operator: adds numbers (2+3=5) OR concatenates strings (\"Hi\"+\"!\" = \"Hi!\")"),
        (0, "Two Types:"),
        (1, "1️⃣  Compile-time Polymorphism — Method Overloading"),
        (1, "2️⃣  Runtime Polymorphism    — Method Overriding"),
    ], slide_num=12)

    # ------------------------------------------------------------------
    # Slide 13 — Method Overloading (Compile-time)
    # ------------------------------------------------------------------
    add_code_slide(prs,
        "Method Overloading — Compile-time Polymorphism",
        [
            "class Calculator {",
            "    // Two integers",
            "    int add(int a, int b) { return a + b; }",
            "",
            "    // Two doubles",
            "    double add(double a, double b) { return a + b; }",
            "",
            "    // Three integers",
            "    int add(int a, int b, int c) { return a + b + c; }",
            "}",
        ],
        "✅  Same method name add(), but different parameters → decided at COMPILE TIME.",
        slide_num=13,
    )

    # ------------------------------------------------------------------
    # Slide 14 — Method Overriding (Runtime)
    # ------------------------------------------------------------------
    add_code_slide(prs,
        "Method Overriding — Runtime Polymorphism",
        [
            "class Animal {",
            "    void sound() { System.out.println(\"Some sound...\"); }",
            "}",
            "",
            "class Cat extends Animal {",
            "    void sound() { System.out.println(\"Meow!\"); }",
            "}",
            "",
            "class Dog extends Animal {",
            "    void sound() { System.out.println(\"Bark!\"); }",
            "}",
            "",
            "// Runtime decision:",
            "Animal a = new Cat();",
            "a.sound();  // Output: Meow!  (decided at RUNTIME)",
        ],
        "✅  Same method name sound(), but behaviour changes based on the actual object at runtime.",
        slide_num=14,
    )

    # ------------------------------------------------------------------
    # Slide 15 — Quick Comparison Table
    # ------------------------------------------------------------------
    add_table_slide(prs,
        "Quick Comparison Table 📋",
        ["Pillar", "Keyword", "Purpose", "Technique"],
        [
            ("Encapsulation", "Hide Data",   "Security & data protection", "private + Getters/Setters"),
            ("Abstraction",   "Hide Logic",  "Simplicity",                 "Abstract class / Interface"),
            ("Inheritance",   "Reuse Code",  "Code reusability",           "extends / implements"),
            ("Polymorphism",  "Many Forms",  "Flexibility",                "Overloading / Overriding"),
        ],
        slide_num=15,
    )

    # ------------------------------------------------------------------
    # Slide 16 — Summary & Memory Trick
    # ------------------------------------------------------------------
    add_content_slide(prs, "Summary & Memory Trick 🧠", [
        (0, "Remember the 4 Pillars with  \"A PIE\" 🥧"),
        (1, "A  →  Abstraction       (hide HOW it works)"),
        (1, "P  →  Polymorphism      (one name, many forms)"),
        (1, "I  →  Inheritance       (reuse from a parent)"),
        (1, "E  →  Encapsulation     (protect your data)"),
        (0, ""),
        (0, "\"OOP is like baking A PIE — you need all 4 ingredients!\""),
        (0, ""),
        (0, "Quick recap:"),
        (1, "Encapsulation  → Capsule / ATM"),
        (1, "Abstraction    → Car dashboard / TV remote"),
        (1, "Inheritance    → Child inherits from parents"),
        (1, "Polymorphism   → Same person, different roles"),
    ], slide_num=16)

    # ------------------------------------------------------------------
    # Slide 17 — Thank You
    # ------------------------------------------------------------------
    add_content_slide(prs, "Thank You! 🙏", [
        (0, "Questions? 🤔"),
        (0, ""),
        (0, "📚 Next Topic:"),
        (1, "Deep dive into each pillar with hands-on coding"),
        (1, "Practice problems & real project examples"),
        (0, ""),
        (0, "\"Code is like humor. When you have to explain it, it's bad.\""),
        (1, "— Cory House"),
    ], slide_num=17)

    return prs


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
if __name__ == "__main__":
    output_file = "OOP_Pillars_Presentation.pptx"
    prs = build_presentation()
    prs.save(output_file)
    print(f"✅ Presentation generated successfully: {output_file}")
